"""Generate the demo PowerPoint deck.

Run:
    python make_ppt.py

Produces: AI_Code_Generator_Demo.pptx
"""
from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Branding
# ---------------------------------------------------------------------------
PURPLE = RGBColor(0xA1, 0x00, 0xFF)
DARK = RGBColor(0x2D, 0x1B, 0x69)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF5, 0xF3, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x10, 0xB9, 0x81)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def set_bg(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    # send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def add_textbox(slide, left, top, width, height, text, *, size=18, bold=False,
                color=DARK, align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, *, size=18, color=DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_accent_bar(slide, top=Inches(1.05)) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), top, Inches(0.6), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE
    bar.line.fill.background()


def add_header(slide, title: str) -> None:
    add_textbox(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.6),
                title, size=30, bold=True, color=DARK)
    add_accent_bar(slide)


def add_footer(slide, page_num: int) -> None:
    add_textbox(slide, Inches(0.6), Inches(7.0), Inches(8), Inches(0.3),
                "AI Code Generator  •  Demo", size=10, color=GREY)
    add_textbox(slide, Inches(12.4), Inches(7.0), Inches(0.6), Inches(0.3),
                str(page_num), size=10, color=GREY)


def add_notes(slide, notes: str) -> None:
    slide.notes_slide.notes_text_frame.text = notes


# ---------------------------------------------------------------------------
# Build deck
# ---------------------------------------------------------------------------
def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    page = 0

    # --- Slide 1: Title -----------------------------------------------------
    page += 1
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK)

    # decorative accent block
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), SLIDE_H)
    accent.fill.solid(); accent.fill.fore_color.rgb = PURPLE; accent.line.fill.background()

    add_textbox(s, Inches(1.0), Inches(2.2), Inches(11), Inches(1.2),
                "AI Code Generator", size=54, bold=True, color=WHITE)
    add_textbox(s, Inches(1.0), Inches(3.3), Inches(11), Inches(0.8),
                "From PRD to Production-Ready Code — Powered by LLMs",
                size=24, color=LIGHT)
    add_textbox(s, Inches(1.0), Inches(5.5), Inches(11), Inches(0.4),
                "Proof of Concept  •  Demo", size=16, color=PURPLE, bold=True)
    add_textbox(s, Inches(1.0), Inches(5.9), Inches(11), Inches(0.4),
                "Ritwika Law", size=14, color=LIGHT)
    add_notes(s,
              "Opening line: 'What if writing software started with one paragraph "
              "and ended with tested, documented code in minutes — not weeks?' "
              "Introduce yourself and the POC.")

    # --- Slide 2: The Problem ---------------------------------------------
    page += 1
    s = prs.slides.add_slide(blank); set_bg(s, WHITE)
    add_header(s, "The Problem")
    add_bullets(s, Inches(0.6), Inches(1.6), Inches(12), Inches(5),
                [
                    "Translating a PRD into a working codebase is slow and repetitive.",
                    "Architecture, code scaffolding and tests are written by hand — even when patterns are well-known.",
                    "Developers spend significant time on boilerplate instead of business logic.",
                    "Onboarding new engineers to a project is expensive — system context lives in heads, not artifacts.",
                    "Tests are often written last (or skipped) under delivery pressure.",
                ], size=20)
    add_footer(s, page)
    add_notes(s,
              "Set the pain point. Most teams know this. Don't oversell — frame it as "
              "a productivity opportunity, not a replacement for engineers.")

    # --- Slide 3: The Idea -------------------------------------------------
    page += 1
    s = prs.slides.add_slide(blank); set_bg(s, WHITE)
    add_header(s, "The Idea")
    add_textbox(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.6),
                "One paragraph in. A full project skeleton out.",
                size=22, bold=True, color=PURPLE)
    add_bullets(s, Inches(0.6), Inches(2.4), Inches(12), Inches(4.2),
                [
                    "Feed a PRD to an LLM-powered pipeline.",
                    "Auto-generate a system design document — architecture, schema, API contracts.",
                    "Auto-generate a production-grade codebase in your chosen stack.",
                    "Auto-generate unit & integration tests against that code.",
                    "Developers review, refine, and ship — instead of starting from scratch.",
                ], size=20)
    add_footer(s, page)
    add_notes(s,
              "Keep this tight. The whole pitch fits on one slide. "
              "Position as a 'first draft generator', not 'auto-pilot'.")

    # --- Slide 4: Pipeline / Architecture ----------------------------------
    page += 1
    s = prs.slides.add_slide(blank); set_bg(s, WHITE)
    add_header(s, "How It Works  —  4-Stage Pipeline")

    stages = [
        ("PRD", "Product Requirements"),
        ("System Design", "Architecture · Schema · APIs"),
        ("Code", "Routes · Services · Models"),
        ("Tests", "Unit · Integration · Edge cases"),
    ]
    box_w = Inches(2.7); box_h = Inches(1.6)
    gap = Inches(0.25)
    total_w = box_w * len(stages) + gap * (len(stages) - 1)
    start_left = (SLIDE_W - total_w) / 2
    top = Inches(2.4)

    for i, (title, sub) in enumerate(stages):
        left = start_left + (box_w + gap) * i
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h)
        shape.fill.solid(); shape.fill.fore_color.rgb = PURPLE if i == 0 else LIGHT
        shape.line.color.rgb = PURPLE
        tf = shape.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        from pptx.enum.text import PP_ALIGN
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = title
        r.font.size = Pt(20); r.font.bold = True
        r.font.color.rgb = WHITE if i == 0 else DARK
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(11)
        r2.font.color.rgb = WHITE if i == 0 else GREY

        if i < len(stages) - 1:
            arrow_left = left + box_w + Inches(0.02)
            arrow = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arrow_left, top + Inches(0.65),
                Inches(0.21), Inches(0.3))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = PURPLE
            arrow.line.fill.background()

    add_textbox(s, Inches(0.6), Inches(4.5), Inches(12), Inches(0.6),
                "Each stage feeds the next. The LLM gets richer context at every step.",
                size=16, color=GREY)
    add_bullets(s, Inches(0.6), Inches(5.1), Inches(12), Inches(2),
                [
                    "Stateless: any stage can be re-run with edits to its input.",
                    "Stack-aware: target Node.js / Express, Java / Spring Boot, or Python / FastAPI.",
                    "Output is plain Markdown + code blocks — easy to copy, diff, and commit.",
                ], size=16)
    add_footer(s, page)
    add_notes(s,
              "Walk through left-to-right. Emphasise that each stage's output becomes the "
              "next stage's input — that's how we keep LLM hallucinations grounded.")

    # --- Slide 5: Tech Stack ----------------------------------------------
    page += 1
    s = prs.slides.add_slide(blank); set_bg(s, WHITE)
    add_header(s, "Tech Stack")

    cols = [
        ("LLM Layer", ["OpenAI / OpenRouter API", "Model: gpt-4o-mini (configurable)", "Temperature: 0.3 for determinism"]),
        ("Backend",   ["Python 3.11+", "Prompt templates per stage", "Pluggable target stacks"]),
        ("UI",        ["Streamlit", "Paste / Upload / Examples", "Live progress + download"]),
        ("Output Stacks", ["Node.js + Express + Jest", "Java + Spring Boot + JUnit", "Python + FastAPI + pytest"]),
    ]
    col_w = Inches(2.95); col_h = Inches(4.5); gap = Inches(0.15)
    total_w = col_w * 4 + gap * 3
    start_left = (SLIDE_W - total_w) / 2
    top = Inches(1.8)
    for i, (title, items) in enumerate(cols):
        left = start_left + (col_w + gap) * i
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, col_w, col_h)
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = PURPLE
        add_textbox(s, left + Inches(0.2), top + Inches(0.2), col_w - Inches(0.4), Inches(0.5),
                    title, size=18, bold=True, color=PURPLE)
        add_bullets(s, left + Inches(0.2), top + Inches(0.8), col_w - Inches(0.4), col_h - Inches(1),
                    items, size=13)
    add_footer(s, page)
    add_notes(s,
              "Keep this fast. The point is: nothing exotic. Standard Python, standard "
              "Streamlit, swap the model with one config line.")

    # --- Slide 6: Live Demo (Screenshots) ----------------------------------
    page += 1
    s = prs.slides.add_slide(blank); set_bg(s, DARK)
    add_textbox(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.6),
                "Live Demo: Step-by-Step", size=30, bold=True, color=WHITE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.05), Inches(0.6), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = PURPLE; bar.line.fill.background()
    add_textbox(s, Inches(1.0), Inches(2.0), Inches(11.0), Inches(1.2),
                "The next slides walk through the application demo using the 'Task Manager' PRD and Node.js stack. Each step is illustrated with real screenshots.", size=20, color=LIGHT)
    add_footer(s, page)
    add_notes(s, "The following slides show the app in action using the 'Task Manager' PRD and Node.js stack.")

    # --- Slide 7: Home/Landing Page ----------------------------------------
    page += 1
    s = prs.slides.add_slide(blank); set_bg(s, WHITE)
    add_header(s, "Step 1: Home/Landing Page")
    img_path = "screenshots/home.png"
    try:
        s.shapes.add_picture(img_path, Inches(2.0), Inches(1.5), width=Inches(9.5))
    except Exception:
        add_textbox(s, Inches(2.0), Inches(3.0), Inches(9.5), Inches(1.0), f"[Missing: {img_path}]", size=18, color=PURPLE)
    add_textbox(s, Inches(2.0), Inches(6.5), Inches(9.5), Inches(0.5), "The app's landing page where users start the code generation process.", size=16, color=GREY)
    add_footer(s, page)

    # --- Slide 8: PRD Input ------------------------------------------------
    page += 1
    s = prs.slides.add_slide(blank); set_bg(s, WHITE)
    add_header(s, "Step 2: Enter PRD Example")
    img_path = "screenshots/prd_input.png"
    try:
        s.shapes.add_picture(img_path, Inches(2.0), Inches(1.5), width=Inches(9.5))
    except Exception:
        add_textbox(s, Inches(2.0), Inches(3.0), Inches(9.5), Inches(1.0), f"[Missing: {img_path}]", size=18, color=PURPLE)
    add_textbox(s, Inches(2.0), Inches(6.5), Inches(9.5), Inches(0.5), "Paste the 'Task Manager' PRD into the input area.", size=16, color=GREY)
    add_footer(s, page)

    # --- Slide 9: Stack Selection ------------------------------------------
    page += 1
    s = prs.slides.add_slide(blank); set_bg(s, WHITE)
    add_header(s, "Step 3: Select Node.js Stack")
    img_path = "screenshots/stack_selection.png"
    try:
        s.shapes.add_picture(img_path, Inches(2.0), Inches(1.5), width=Inches(9.5))
    except Exception:
        add_textbox(s, Inches(2.0), Inches(3.0), Inches(9.5), Inches(1.0), f"[Missing: {img_path}]", size=18, color=PURPLE)
    add_textbox(s, Inches(2.0), Inches(6.5), Inches(9.5), Inches(0.5), "Choose Node.js + Express as the target stack.", size=16, color=GREY)
    add_footer(s, page)

    # --- Slide 10: Generation/Progress -------------------------------------
    page += 1
    s = prs.slides.add_slide(blank); set_bg(s, WHITE)
    add_header(s, "Step 4: Generation in Progress")
    img_path = "screenshots/generation.png"
    try:
        s.shapes.add_picture(img_path, Inches(2.0), Inches(1.5), width=Inches(9.5))
    except Exception:
        add_textbox(s, Inches(2.0), Inches(3.0), Inches(9.5), Inches(1.0), f"[Missing: {img_path}]", size=18, color=PURPLE)
    add_textbox(s, Inches(2.0), Inches(6.5), Inches(9.5), Inches(0.5), "The app shows real-time progress as the LLM generates code.", size=16, color=GREY)
    add_footer(s, page)

    # --- Slide 11: Output/Results (Three Screenshots) ----------------------
    page += 1
    s = prs.slides.add_slide(blank); set_bg(s, WHITE)
    add_header(s, "Step 5: Output & Download")
    output_names = ["system_design", "code", "tests"]
    img_width = Inches(3.0)
    img_height = Inches(2.5)
    gap = Inches(0.5)
    left_start = (SLIDE_W - (img_width * 3 + gap * 2)) / 2
    top = Inches(2.0)
    captions = ["System Design Output", "Generated Code Output", "Test Cases Output"]
    for i, name in enumerate(output_names):
        img_path = f"screenshots/output_{name}.png"
        left = left_start + i * (img_width + gap)
        try:
            s.shapes.add_picture(img_path, left, top, width=img_width, height=img_height)
        except Exception:
            add_textbox(s, left, top + img_height/2, img_width, Inches(0.5), f"[Missing: {img_path}]", size=14, color=PURPLE)
        add_textbox(s, left, top + img_height + Inches(0.1), img_width, Inches(0.4), captions[i], size=13, color=GREY)
    add_textbox(s, Inches(1.0), Inches(6.5), Inches(11), Inches(0.5), "View and download the generated code, system design, and test cases.", size=16, color=GREY)
    add_footer(s, page)

    # --- Slide 7: Value & Benefits ----------------------------------------
    page += 1
    s = prs.slides.add_slide(blank); set_bg(s, WHITE)
    add_header(s, "Why It Matters")

    benefits = [
        ("Speed",        "Hours of scaffolding compressed into minutes."),
        ("Consistency",  "Every project starts from the same architectural baseline."),
        ("Documentation","System design doc is a by-product, not an afterthought."),
        ("Test-first",   "Tests are generated alongside code — not skipped under deadline."),
        ("Onboarding",   "New devs read a clean, generated baseline before customising."),
        ("Stack-agnostic","Same PRD, different language — useful for prototyping and migration."),
    ]
    col_w = Inches(6.0); row_h = Inches(1.3); gap_x = Inches(0.3); gap_y = Inches(0.15)
    top0 = Inches(1.7); left0 = Inches(0.6)
    for i, (title, desc) in enumerate(benefits):
        row, col = divmod(i, 2)
        left = left0 + (col_w + gap_x) * col
        top = top0 + (row_h + gap_y) * row
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, col_w, row_h)
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = PURPLE
        add_textbox(s, left + Inches(0.25), top + Inches(0.15), col_w - Inches(0.5), Inches(0.45),
                    title, size=18, bold=True, color=PURPLE)
        add_textbox(s, left + Inches(0.25), top + Inches(0.6), col_w - Inches(0.5), Inches(0.65),
                    desc, size=14, color=DARK)
    add_footer(s, page)
    add_notes(s,
              "Lead with speed — it's the easiest win to communicate. "
              "End with 'this is a starting point, not a replacement'.")

    # --- Slide 8: Limitations & Honest View ---------------------------------
    page += 1
    s = prs.slides.add_slide(blank); set_bg(s, WHITE)
    add_header(s, "What This Is — and What It Isn't")
    add_textbox(s, Inches(0.6), Inches(1.6), Inches(6), Inches(0.5),
                "Strengths", size=22, bold=True, color=ACCENT)
    add_bullets(s, Inches(0.6), Inches(2.1), Inches(6), Inches(4.5),
                [
                    "Excellent for greenfield prototypes.",
                    "Strong on well-known patterns (REST CRUD, auth, etc.).",
                    "Produces readable, idiomatic code.",
                    "Tests cover happy + key edge paths.",
                ], size=16)
    add_textbox(s, Inches(7.0), Inches(1.6), Inches(6), Inches(0.5),
                "Honest Limitations", size=22, bold=True, color=PURPLE)
    add_bullets(s, Inches(7.0), Inches(2.1), Inches(6), Inches(4.5),
                [
                    "Not a replacement for an engineer — it's a first-draft generator.",
                    "Generated code needs review for security & correctness.",
                    "Complex domain logic still needs human design.",
                    "LLM cost scales with PRD size and number of regenerations.",
                ], size=16)
    add_footer(s, page)
    add_notes(s,
              "Showing limitations builds credibility. Audience trusts you more when "
              "you name the trade-offs before they do.")

    # --- Slide 9: Roadmap -------------------------------------------------
    page += 1
    s = prs.slides.add_slide(blank); set_bg(s, WHITE)
    add_header(s, "What's Next")
    roadmap = [
        ("Near-term",  ["Multi-file output as a runnable repo (not just markdown).",
                        "Git-init + first-commit automation.",
                        "Diagram generation (Mermaid) alongside design doc."]),
        ("Mid-term",   ["Agentic refinement loop — generator critiques and self-edits.",
                        "Plug in retrieval over org-specific patterns and guardrails.",
                        "CI pipeline scaffolding (GitHub Actions / Azure DevOps)."]),
        ("Long-term",  ["Iterative dev loop: developer edits → LLM re-syncs design & tests.",
                        "Compliance / security checks built into the generation pipeline.",
                        "Multi-service architecture generation from a single PRD."]),
    ]
    col_w = Inches(4.15); col_h = Inches(5.0); gap = Inches(0.15)
    total_w = col_w * 3 + gap * 2
    start_left = (SLIDE_W - total_w) / 2
    top = Inches(1.7)
    for i, (title, items) in enumerate(roadmap):
        left = start_left + (col_w + gap) * i
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, col_w, col_h)
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = PURPLE
        add_textbox(s, left + Inches(0.25), top + Inches(0.2), col_w - Inches(0.5), Inches(0.5),
                    title, size=20, bold=True, color=PURPLE)
        add_bullets(s, left + Inches(0.25), top + Inches(0.9), col_w - Inches(0.5), col_h - Inches(1),
                    items, size=14)
    add_footer(s, page)
    add_notes(s,
              "Pitch this as a roadmap, not a wishlist. Invite collaborators if your "
              "forum is the right audience.")

    # --- Slide 10: Thank You / Q&A ----------------------------------------
    page += 1
    s = prs.slides.add_slide(blank); set_bg(s, DARK)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), SLIDE_H)
    accent.fill.solid(); accent.fill.fore_color.rgb = PURPLE; accent.line.fill.background()
    add_textbox(s, Inches(1.0), Inches(2.8), Inches(11), Inches(1.5),
                "Thank you.", size=64, bold=True, color=WHITE)
    add_textbox(s, Inches(1.0), Inches(4.0), Inches(11), Inches(0.8),
                "Questions & feedback welcome.", size=24, color=LIGHT)
    add_textbox(s, Inches(1.0), Inches(6.6), Inches(11), Inches(0.4),
                "ritwika.law@accenture.com", size=14, color=PURPLE)
    add_notes(s,
              "Close with: 'I'd love to hear which parts of your team's workflow this "
              "could plug into.' Open the floor.")

    return prs


def main() -> None:
    prs = build()
    out = "AI_Code_Generator_Demo.pptx"
    prs.save(out)
    print(f"Wrote {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
